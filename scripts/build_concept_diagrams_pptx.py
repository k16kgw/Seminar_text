"""教材の概念図に対応する編集可能なPowerPointを生成する．

各要素はPowerPointの文字，図形，矢印として作る．本文表示用のSVGを
貼り付けたものではないため，PowerPoint上で色，位置，文言を変更できる．
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import subprocess
import tempfile
from zipfile import ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "assets" / "figures" / "concepts" / "concept_diagrams_editable.pptx"

NAVY = "173451"
TEXT = "21364A"
MUTED = "52677A"
LINE = "60798E"
BLUE = "EAF4FF"
GREEN = "EAF8EF"
ORANGE = "FFF3E1"
PURPLE = "F3EDFF"
RED = "FFECEF"
TEAL = "EAF8F7"
WHITE = "FFFFFF"
BG = "FBFDFF"

A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
MATH: dict[str, etree._Element] = {}

FORMULAS = {
    "stencil": r"\Delta u_{i,j}=\frac{u_{i+1,j}+u_{i-1,j}+u_{i,j+1}+u_{i,j-1}-4u_{i,j}}{(\Delta x)^2}",
    "delta_negative": r"\Delta u_{i,j}<0",
    "delta_positive": r"\Delta u_{i,j}>0",
    "reaction": r"u\rightleftarrows v",
    "adjacency": r"A=\begin{pmatrix}0&0.8&0\\0&0&1.0\\0.5&0&0\end{pmatrix}",
    "matrix_action": r"(A\mathbf{x})_1=0.8x_2",
    "base_temperature": r"T=T_b",
    "fin_conduction": r"-kA\frac{dT}{dx}",
    "fin_convection": r"hP(T-T_\infty)",
    "plate_equation": r"k\Delta T",
    "plate_root": r"T=T_b",
    "plate_boundary": r"-k\frac{\partial T}{\partial n}=h(T-T_\infty)",
    "seed_runs": r"\mathrm{seed}\ 1\mapsto T_1,\ldots,\mathrm{seed}\ n\mapsto T_n",
    "state_r": r"R(t)",
    "state_j": r"J(t)",
    "two_person_ode": r"\frac{d}{dt}\begin{pmatrix}R\\J\end{pmatrix}=M\begin{pmatrix}R\\J\end{pmatrix}",
    "two_person_matrix": r"M=\begin{pmatrix}a&b\\c&d\end{pmatrix}",
    "coefficient_b": r"b",
    "coefficient_c": r"c",
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_text(shape, text: str, size: int = 18, bold: bool = False,
             color: str = TEXT, align: PP_ALIGN = PP_ALIGN.CENTER) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(slide, text: str, x: float, y: float, w: float, h: float,
             size: int = 18, bold: bool = False, color: str = TEXT,
             align: PP_ALIGN = PP_ALIGN.CENTER):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, bold, color, align)
    return shape


def add_box(slide, text: str, x: float, y: float, w: float, h: float,
            fill: str = BLUE, size: int = 18, bold: bool = False,
            line: str = LINE, radius: bool = True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.5)
    set_text(shape, text, size, bold)
    return shape


def add_circle(slide, text: str, x: float, y: float, d: float,
               fill: str = BLUE, line: str = LINE, size: int = 19):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(1.8)
    set_text(shape, text, size, True)
    return shape


def add_arrow(slide, x: float, y: float, w: float, h: float, direction: str = "right",
              fill: str = LINE):
    kinds = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    shape = slide.shapes.add_shape(kinds[direction], Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(fill)
    return shape


def build_math_catalog() -> dict[str, etree._Element]:
    """LaTeX数式をPandocでOffice Mathへ変換する．"""
    markdown = "## 数式\n\n" + "\n\n".join(
        f"$${formula}$$" for formula in FORMULAS.values()
    )
    with tempfile.TemporaryDirectory() as tmp:
        tmpdir = Path(tmp)
        source = tmpdir / "equations.md"
        output = tmpdir / "equations.pptx"
        source.write_text(markdown, encoding="utf-8")
        subprocess.run(
            ["pandoc", str(source), "-o", str(output), "--slide-level=2"],
            check=True,
            capture_output=True,
            text=True,
        )
        with ZipFile(output) as archive:
            xml = archive.read("ppt/slides/slide1.xml")
    root = etree.fromstring(xml)
    elements = root.xpath(".//a14:m", namespaces={"a14": A14_NS})
    if len(elements) != len(FORMULAS):
        raise RuntimeError(
            f"Office Math変換数が一致しない: expected={len(FORMULAS)}, actual={len(elements)}"
        )
    return {
        key: deepcopy(element)
        for key, element in zip(FORMULAS, elements)
    }


def add_math(slide, key: str, x: float, y: float, w: float, h: float,
             size: int = 20):
    """編集可能なOffice Math数式をテキストボックスへ追加する．"""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]._p
    for child in list(paragraph):
        paragraph.remove(child)
    p_pr = OxmlElement("a:pPr")
    p_pr.set("algn", "ctr")
    default_run = OxmlElement("a:defRPr")
    default_run.set("sz", str(size * 100))
    default_run.set("lang", "ja-JP")
    p_pr.append(default_run)
    paragraph.append(p_pr)
    paragraph.append(deepcopy(MATH[key]))
    end_run = OxmlElement("a:endParaRPr")
    end_run.set("sz", str(size * 100))
    end_run.set("lang", "ja-JP")
    paragraph.append(end_run)
    return shape


def base_slide(prs: Presentation, title: str, source: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(BG)
    add_text(slide, title, 0.35, 0.18, 12.63, 0.5, 24, True, NAVY)
    add_text(slide, f"対応SVG: {source}", 0.35, 7.12, 12.63, 0.22, 9, False, MUTED, PP_ALIGN.RIGHT)
    return slide


def pipeline(slide, items: list[tuple[str, str]], y: float = 2.2,
             note: str | None = None) -> None:
    colors = [BLUE, GREEN, ORANGE, PURPLE, TEAL, RED]
    n = len(items)
    gap = 0.18
    margin = 0.35
    arrow_w = 0.28
    total_arrow = arrow_w * (n - 1)
    box_w = (13.333 - 2 * margin - total_arrow - gap * 2 * (n - 1)) / n
    x = margin
    for i, (heading, detail) in enumerate(items):
        add_box(slide, f"{heading}\n{detail}", x, y, box_w, 1.45, colors[i % len(colors)], 15, i == 0)
        x += box_w
        if i < n - 1:
            add_arrow(slide, x + gap, y + 0.57, arrow_w, 0.31)
            x += arrow_w + 2 * gap
    if note:
        add_text(slide, note, 0.7, y + 1.8, 11.93, 0.55, 16, False, MUTED)


def cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)
    add_text(slide, "卒業研究準備セミナー\n概念図・編集用元データ", 1.2, 1.45, 10.93, 1.4, 30, True, NAVY)
    add_text(slide, "すべての文字・図形・矢印はPowerPoint上で個別に編集できます．", 1.2, 3.35, 10.93, 0.65, 18, False, TEXT)
    add_text(slide, "本文表示用SVGと同じ考え方を，編集しやすい図形で再構成しています．", 1.2, 4.15, 10.93, 0.65, 16, False, MUTED)


def learning_path(prs: Presentation) -> None:
    slide = base_slide(prs, "基礎学習から卒業研究へ進む", "learning_path.svg")
    pipeline(slide, [
        ("卒研の準備", "研究設計・各モデルの入口"),
        ("基礎モデル", "最小の式・標準設定"),
        ("研究技能", "測る・比較する・検証する"),
        ("ロードマップ", "必要な段階を参照する"),
        ("卒業研究", "拡張・比較・考察"),
    ], 2.2, "問題が見つかったら，基礎モデルや研究技能へ戻る．")


def research_cycle(prs: Presentation) -> None:
    slide = base_slide(prs, "数理モデリングは検証しながら戻る循環", "research_cycle.svg")
    top = [("現象と問い", "何を比較するか"), ("状態と仮定", "何を残すか"),
           ("最小モデル", "式・更新規則"), ("数値実験", "条件を1つ変える")]
    pipeline(slide, top, 1.45)
    bottom = [("妥当性と限界", "どこまで言えるか"), ("検証", "誤差・再現性"),
              ("評価指標", "図を数値へ変える")]
    pipeline(slide, bottom, 4.25, "結果が問いに答えなければ，状態変数・仮定・比較条件を見直す．")
    add_arrow(slide, 11.8, 3.05, 0.35, 0.7, "down")
    add_arrow(slide, 1.05, 3.05, 0.35, 0.7, "up")


def pde_stencil(prs: Presentation) -> None:
    slide = base_slide(prs, "2次元ラプラシアンの5点差分", "pde_five_point_stencil.svg")
    x0, y0, cell = 0.65, 1.15, 0.78
    for row in range(5):
        for col in range(5):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0 + col * cell), Inches(y0 + row * cell), Inches(cell), Inches(cell))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(WHITE)
            shape.line.color.rgb = rgb("A8B8C5")
    points = {(2, 2): ("中央", "FFDB8A"), (1, 2): ("上", "DCEEFF"),
              (3, 2): ("下", "DCEEFF"), (2, 1): ("左", "DCEEFF"),
              (2, 3): ("右", "DCEEFF")}
    for (row, col), (label, color) in points.items():
        add_circle(slide, label, x0 + col * cell + 0.12, y0 + row * cell + 0.12, 0.54, color, "3979AD", 13)
    add_box(slide, "", 5.25, 1.55, 7.1, 1.4, BLUE)
    add_math(slide, "stencil", 5.45, 1.7, 6.7, 1.05, 18)
    add_box(slide, "中央が周囲より高い\n\n次時刻に中央の値が下がる", 5.25, 3.45, 3.35, 1.45, ORANGE, 16)
    add_math(slide, "delta_negative", 5.75, 3.95, 2.35, 0.38, 17)
    add_box(slide, "中央が周囲より低い\n\n次時刻に中央の値が上がる", 9.0, 3.45, 3.35, 1.45, GREEN, 16)
    add_math(slide, "delta_positive", 9.5, 3.95, 2.35, 0.38, 17)
    add_text(slide, "中央と上下左右の4点を参照する．斜め4点は含めない．", 5.25, 5.15, 7.1, 0.55, 17, False, MUTED)


def reaction_diffusion(prs: Presentation) -> None:
    slide = base_slide(prs, "反応拡散 = 各点の反応 + 点どうしの拡散", "reaction_diffusion_parts.svg")
    add_box(slide, "反応だけ\n各場所で2変数ODEが進む", 0.6, 1.55, 3.35, 2.3, ORANGE, 19, True)
    add_math(slide, "reaction", 1.45, 3.05, 1.65, 0.42, 19)
    add_text(slide, "＋", 4.05, 2.2, 0.55, 0.6, 28, True, LINE)
    add_box(slide, "拡散だけ\n高い場所から低い場所へ広がる\n空間差をならす", 4.7, 1.55, 3.35, 2.3, BLUE, 19, True)
    add_arrow(slide, 8.25, 2.45, 0.6, 0.35)
    add_box(slide, "組合せ\n特定の空間スケールが成長\n斑点・縞・迷路", 9.05, 1.55, 3.65, 2.3, PURPLE, 19, True)
    add_text(slide, "式を読むときは，まず反応だけ・拡散だけの場合を分けて考える．", 1.0, 4.55, 11.33, 0.65, 18, False, MUTED)


def agent_updates(prs: Presentation) -> None:
    slide = base_slide(prs, "逐次更新と一斉更新の違い", "agent_update_modes.svg")
    add_box(slide, "更新前\nA → □ ← B\n2人が中央セルを希望", 0.55, 1.45, 3.55, 2.4, BLUE, 19, True)
    add_arrow(slide, 4.25, 2.45, 0.55, 0.35)
    add_box(slide, "逐次更新\nAを先に処理\nAが中央を確保する", 4.95, 1.45, 3.25, 2.4, GREEN, 19, True)
    add_arrow(slide, 8.35, 2.45, 0.55, 0.35)
    add_box(slide, "一斉更新\n全員の希望を集める\n競合規則で結果を決める", 9.05, 1.45, 3.7, 2.4, PURPLE, 19, True)
    add_text(slide, "更新順序は実装上の都合ではなく，明記すべきモデルの仮定である．", 0.8, 4.65, 11.73, 0.7, 18, False, MUTED)


def network_matrix(prs: Presentation) -> None:
    slide = base_slide(prs, "有向グラフと隣接行列の対応", "network_matrix_graph.svg")
    add_circle(slide, "1", 1.95, 1.25, 1.0, BLUE)
    add_circle(slide, "2", 0.75, 3.8, 1.0, GREEN)
    add_circle(slide, "3", 3.15, 3.8, 1.0, ORANGE)
    add_box(slide, "2 → 1 : 0.8\n3 → 2 : 1.0\n1 → 3 : 0.5", 1.25, 2.5, 2.5, 1.0, WHITE, 16)
    add_box(slide, "", 5.0, 1.45, 6.9, 2.25, BLUE)
    add_math(slide, "adjacency", 5.65, 1.7, 5.6, 1.7, 21)
    add_box(slide, "第1行は頂点1へ入る影響", 5.75, 4.2, 5.4, 1.35, GREEN, 17)
    add_math(slide, "matrix_action", 7.05, 4.32, 2.8, 0.48, 18)
    add_text(slide, "Aᵢⱼ は頂点 j から頂点 i への影響．列が出発点，行が到着点．", 0.8, 6.0, 11.73, 0.55, 17, False, MUTED)


def fin_balance(prs: Presentation) -> None:
    slide = base_slide(prs, "1次元フィンの熱収支", "fin_heat_balance.svg")
    add_box(slide, "体部\n根元温度を固定", 0.55, 2.0, 1.4, 2.1, "E85F49", 17, True)
    add_math(slide, "base_temperature", 0.68, 3.22, 1.15, 0.38, 16)
    add_box(slide, "フィン内部の熱伝導", 2.0, 2.3, 9.5, 1.45, "FFF0D0", 22, True, "A17839", False)
    add_math(slide, "fin_conduction", 5.25, 3.02, 3.0, 0.48, 19)
    add_arrow(slide, 1.55, 2.8, 1.0, 0.35, "right", "D65D47")
    for x in [3.1, 5.0, 6.9, 8.8, 10.7]:
        add_arrow(slide, x, 1.35, 0.35, 0.75, "up")
        add_arrow(slide, x, 3.95, 0.35, 0.75, "down")
    add_text(slide, "表面から周囲へ対流で放熱する", 3.0, 0.78, 7.8, 0.42, 18, True, TEXT)
    add_math(slide, "fin_convection", 5.1, 1.16, 3.6, 0.42, 18)
    add_text(slide, "先へ進むほど熱が周囲へ失われ，温度差が小さくなる．", 1.5, 5.35, 10.4, 0.65, 18, False, MUTED)


def plate_boundaries(prs: Presentation) -> None:
    slide = base_slide(prs, "薄板モデルの領域と境界条件", "plate_boundary_conditions.svg")
    add_box(slide, "板の内部領域\n熱伝導方程式を解く", 1.0, 1.2, 6.2, 3.8, "FFF4DD", 23, True, "B98A35", False)
    add_math(slide, "plate_equation", 3.25, 3.15, 1.7, 0.5, 20)
    add_box(slide, "根元は固定温度", 1.0, 4.45, 6.2, 0.65, "EE6453", 18, True, "EE6453", False)
    add_math(slide, "plate_root", 5.3, 4.55, 1.35, 0.38, 17)
    for x in [1.35, 2.65, 3.95, 5.25, 6.55]:
        add_arrow(slide, x, 0.75, 0.32, 0.55, "up")
    add_arrow(slide, 0.45, 2.3, 0.55, 0.32, "left")
    add_arrow(slide, 7.2, 2.3, 0.55, 0.32, "right")
    add_box(slide, "計算で区別するもの\n\n内部セル\n根元セル\n対流境界\n形状の外部", 8.25, 1.2, 3.8, 3.8, BLUE, 18, True)
    add_text(slide, "側面・先端の対流境界条件", 1.0, 5.3, 6.8, 0.42, 17, False, MUTED)
    add_math(slide, "plate_boundary", 1.65, 5.72, 5.5, 0.55, 17)


def pattern_metrics(prs: Presentation) -> None:
    slide = base_slide(prs, "模様の見た目を複数の観察量へ分ける", "pattern_measurements.svg")
    items = [("入力する模様", "同じ画像を使う"), ("面積比", "明るい画素 ÷ 全画素"),
             ("個数", "連結領域を数える"), ("特徴波長", "模様の代表的な間隔"),
             ("方向", "縞の向き")]
    pipeline(slide, items, 2.0, "1つの特徴量だけでは区別できない模様がある．")


def boarding_rules(prs: Presentation) -> None:
    slide = base_slide(prs, "乗降セルオートマトンの空間と規則", "boarding_space_rules.svg")
    add_box(slide, "車内", 0.65, 1.1, 7.2, 2.5, BLUE, 22, True, "3979AD", False)
    add_box(slide, "ホーム", 0.65, 3.6, 7.2, 1.4, ORANGE, 22, True, "BD8722", False)
    add_box(slide, "ドア", 3.55, 3.2, 1.4, 1.8, WHITE, 18, True, "D25E4F", False)
    add_circle(slide, "降", 3.65, 1.65, 0.75, RED, "D95F65", 16)
    add_circle(slide, "乗", 4.15, 4.0, 0.75, BLUE, "3C83C5", 16)
    add_arrow(slide, 3.82, 2.4, 0.35, 0.7, "down")
    add_arrow(slide, 4.35, 3.35, 0.35, 0.7, "up")
    add_box(slide, "1ステップの規則\n\n1. 近傍から候補を選ぶ\n2. 同じセルの競合を解く\n3. 確定した移動を反映\n\n1セルに1人", 8.45, 1.25, 4.1, 3.9, PURPLE, 17, True)


def repeated_simulation(prs: Presentation) -> None:
    slide = base_slide(prs, "確率的シミュレーションは分布で比較する", "repeated_simulation.svg")
    pipeline(slide, [("固定する条件", "人数・空間・規則・更新方式"),
                     ("seedだけ変える", "複数のseedで独立に実行する"),
                     ("分布として比較", "平均・中央値・ばらつき・未完了率")],
             2.05, "基準条件と代替条件には同じseed集合を使い，偶然差をそろえる．")
    add_math(slide, "seed_runs", 4.45, 4.45, 4.5, 0.5, 17)


def two_person(prs: Presentation) -> None:
    slide = base_slide(prs, "2人の線形ODEと係数行列", "two_person_coupling.svg")
    add_circle(slide, "", 1.25, 2.05, 1.5, BLUE, "397BB2", 21)
    add_math(slide, "state_r", 1.55, 2.5, 0.9, 0.5, 19)
    add_circle(slide, "", 4.25, 2.05, 1.5, ORANGE, "C86B42", 21)
    add_math(slide, "state_j", 4.55, 2.5, 0.9, 0.5, 19)
    add_arrow(slide, 2.8, 1.75, 1.4, 0.35, "right")
    add_math(slide, "coefficient_c", 3.25, 1.28, 0.5, 0.42, 19)
    add_arrow(slide, 2.8, 3.6, 1.4, 0.35, "left")
    add_math(slide, "coefficient_b", 3.25, 4.0, 0.5, 0.42, 19)
    add_box(slide, "", 7.0, 1.2, 5.1, 4.2, PURPLE)
    add_math(slide, "two_person_ode", 7.35, 1.6, 4.4, 1.0, 18)
    add_math(slide, "two_person_matrix", 8.1, 2.85, 2.9, 1.05, 20)
    add_text(slide, "第1行はRの変化率\n第2行はJの変化率", 7.55, 4.15, 4.0, 0.75, 16, False, MUTED)
    add_text(slide, "符号を1個ずつ見るだけでなく，行列全体の固有値で挙動を判断する．", 1.0, 5.75, 11.33, 0.6, 17, False, MUTED)


def theme_pipelines(prs: Presentation) -> None:
    data = [
        ("背板形状比較を研究へ育てる流れ", "heat_research_pipeline.svg",
         [("基準計算", "長方形・収支・格子"), ("形状を1つ追加", "輪郭・面積"),
          ("公平に正規化", "同体積など"), ("2形状を比較", "同じ色尺度・主指標"),
          ("差を検証", "数値誤差との大小"), ("拡張", "形態系列・感度")]),
        ("実画像と生成模様を比較する流れ", "pattern_research_pipeline.svg",
         [("実画像", "前処理を固定"), ("生成模様", "seed・保存時刻"),
          ("同じ特徴量", "面積比・波長など"), ("距離を比較", "粗い探索"),
          ("検証", "未使用画像・別seed")]),
        ("乗降行動を規則へ変換して比較する流れ", "boarding_research_pipeline.svg",
         [("観察", "位置変化を記録"), ("操作的定義", "いつ・どこへ"),
          ("擬似コード", "競合も記述"), ("小配置テスト", "1～3人"),
          ("反復比較", "同じseed集合"), ("拡張", "割合・人数・感度")]),
        ("物語の観測から関係モデルを検証する流れ", "relation_research_pipeline.svg",
         [("対象範囲", "人物・場面・時間"), ("状態を定義", "尺度・不明値"),
          ("少量で試す", "不一致から修正"), ("2人基準", "予測と残差"),
          ("1項だけ追加", "残差に対応"), ("検証", "未使用場面")]),
    ]
    for title, source, items in data:
        slide = base_slide(prs, title, source)
        pipeline(slide, items, 2.15, "一度に複雑化せず，各段階の判断と検証結果を研究ノートに残す．")


def cross_theme(prs: Presentation) -> None:
    slide = base_slide(prs, "4テーマから研究方法を相互に借りる", "cross_theme_transfer.svg")
    add_box(slide, "自分の卒業研究\n問い → 最小モデル → 比較\n→ 評価指標 → 検証 → 限界", 4.35, 2.35, 4.65, 2.05, "EDF3F8", 20, True)
    add_box(slide, "放熱\n公平な比較条件\n格子収束・収支", 0.6, 1.0, 3.1, 1.65, RED, 18, True)
    add_box(slide, "模様\n見た目を特徴量へ\nパラメータ地図", 9.65, 1.0, 3.1, 1.65, PURPLE, 18, True)
    add_box(slide, "乗降\n更新規則への感度\nseed反復・分布", 0.6, 4.6, 3.1, 1.65, BLUE, 18, True)
    add_box(slide, "関係モデル\n結合構造・固有値\n推定と検証の分離", 9.65, 4.6, 3.1, 1.65, GREEN, 18, True)
    add_arrow(slide, 3.7, 1.8, 0.7, 0.35)
    add_arrow(slide, 8.95, 1.8, 0.7, 0.35, "left")
    add_arrow(slide, 3.7, 5.3, 0.7, 0.35)
    add_arrow(slide, 8.95, 5.3, 0.7, 0.35, "left")


def main() -> None:
    global MATH
    MATH = build_math_catalog()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    learning_path(prs)
    research_cycle(prs)
    pde_stencil(prs)
    reaction_diffusion(prs)
    agent_updates(prs)
    network_matrix(prs)
    fin_balance(prs)
    plate_boundaries(prs)
    pattern_metrics(prs)
    boarding_rules(prs)
    repeated_simulation(prs)
    two_person(prs)
    theme_pipelines(prs)
    cross_theme(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT.relative_to(ROOT)} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
